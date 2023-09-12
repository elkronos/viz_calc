import numpy as np
import pandas as pd
from pptx import Presentation
from io import BytesIO
import matplotlib.pyplot as plt
from pptx.util import Inches
from itertools import combinations
import seaborn as sns

def validate_inputs(data, unique_levels, plots_at_a_time):
    if not isinstance(data, pd.core.frame.DataFrame):
        raise ValueError("The input data is not a valid pandas DataFrame.")
    
    if not isinstance(unique_levels, int) or unique_levels <= 0:
        raise ValueError("unique_levels should be a positive integer.")
    
    if plots_at_a_time != 'all' and (not isinstance(plots_at_a_time, int) or plots_at_a_time <= 0):
        raise ValueError("plots_at_a_time should be either 'all' or a positive integer.")

def get_cols_to_plot(data, unique_levels):
    return [col for col in data.columns if data[col].nunique() >= unique_levels]

def create_scatter_plot(x, y, ax, line_of_best_fit, confidence_interval):
    if line_of_best_fit:
        sns.regplot(x=x, y=y, ax=ax, ci=95 if confidence_interval else None, scatter_kws={'s':10})
    else:
        ax.scatter(x, y)

def set_plot_properties(ax, col1, col2):
    ax.set_ylabel(col2)
    ax.set_xlabel(col1)
    ax.set_title(f'Scatter Plot of {col1} and {col2}')

def all_scatters(data: pd.core.frame.DataFrame, unique_levels: int = 10, save_to_ppt: bool = False, line_of_best_fit: bool = False, confidence_interval: bool = False, ppt_filename: str = 'output.pptx', plots_at_a_time: int = 'all') -> None:
    try:
        validate_inputs(data, unique_levels, plots_at_a_time)
    except ValueError as e:
        print(e)
        return

    cols_to_plot = get_cols_to_plot(data, unique_levels)

    if not cols_to_plot:
        print("No columns to plot.")
        return

    if save_to_ppt:
        prs = Presentation()

    column_combinations = list(combinations(cols_to_plot, 2))
    num_plots = len(column_combinations)
    if plots_at_a_time == 'all':
        plots_at_a_time = num_plots

    for i, (col1, col2) in enumerate(column_combinations):
        x = data[col1]
        y = data[col2]
        
        fig, ax = plt.subplots()
        
        create_scatter_plot(x, y, ax, line_of_best_fit, confidence_interval)
        set_plot_properties(ax, col1, col2)

        if save_to_ppt:
            img_stream = BytesIO()
            plt.tight_layout()
            plt.savefig(img_stream, format='png')
            plt.close()
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            left = top = Inches(1)
            slide.shapes.add_picture(img_stream, left, top, Inches(4.5), Inches(4.5))
        else:
            plt.tight_layout()
            plt.show()
        
        if (i+1) % plots_at_a_time == 0 and i < num_plots - 1:
            input("Press Enter to see next set of plots...")
            
    if save_to_ppt:
        prs.save(ppt_filename)
        print(f'Saved PowerPoint file to: {ppt_filename}')

# Usage example
np.random.seed(123)
data = {
    'var1': np.random.choice(list(range(10)), 100),  # < 10 unique levels
    'var2': np.random.choice(list(range(10)), 100),  # < 10 unique levels
    'var3': np.random.choice(list(range(10)), 100),  # < 10 unique levels
    'var4': np.random.choice(list(range(20)), 100),  # 10-20 unique levels
    'var5': np.random.choice(list(range(20)), 100),  # 10-20 unique levels
    'var6': np.random.choice(list(range(20)), 100),  # 10-20 unique levels
    'var7': np.random.choice(list(range(200)), 100), # > 100 unique levels
    'var8': np.random.choice(list(range(200)), 100), # > 100 unique levels
    'var9': np.random.choice(list(range(200)), 100), # > 100 unique levels
    'var10': np.random.choice(list(range(200)), 100) # > 100 unique levels
}
df = pd.DataFrame(data)
all_scatters(df, unique_levels=15, save_to_ppt=False, line_of_best_fit=True, confidence_interval=True, plots_at_a_time=4)
