import numpy as np
import pandas as pd
from pptx import Presentation
from io import BytesIO
import matplotlib.pyplot as plt
from pptx.util import Inches
from itertools import combinations
import seaborn as sns


class InvalidInputError(ValueError):
    pass


def validate_inputs(data, categorical_min, continuous_max, plots_at_a_time):
    """Validates the inputs for the all_boxes function."""
    if not isinstance(data, pd.core.frame.DataFrame):
        raise InvalidInputError("The input data is not a valid pandas DataFrame.")

    if not isinstance(categorical_min, int) or categorical_min <= 1:
        raise InvalidInputError("categorical_min should be an integer greater than 1.")

    if not isinstance(continuous_max, int) or continuous_max <= 0:
        raise InvalidInputError("continuous_max should be a positive integer.")

    if plots_at_a_time != 'all' and (not isinstance(plots_at_a_time, int) or plots_at_a_time <= 0):
        raise InvalidInputError("plots_at_a_time should be either 'all' or a positive integer.")


def get_cols_to_plot(data, categorical_min, continuous_max):
    """Returns columns satisfying categorical_min and continuous_max conditions."""
    return [col for col in data.columns if data[col].nunique() >= categorical_min and data[col].max() <= continuous_max]


def set_plot_properties(ax, col1, col2):
    """Sets the properties for the plot."""
    ax.set_ylabel(col2)
    ax.set_xlabel(col1)
    ax.set_title(f'Box Plot of {col1} and {col2}')


def save_plot_to_ppt(prs, img_stream):
    """Saves the plot to PowerPoint presentation."""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    left = top = Inches(1)
    slide.shapes.add_picture(img_stream, left, top, Inches(4.5), Inches(4.5))


def all_boxes(data, categorical_min=2, continuous_max=100, save_to_ppt=False, ppt_filename='output.pptx', plots_at_a_time='all'):
    """
    Generates box plots for combinations of columns in the input data frame based on the specified conditions.
    
    The function plots box plots for each combination of two columns in the data frame which satisfy the conditions 
    defined by `categorical_min` and `continuous_max`. These plots can either be displayed or saved to a PowerPoint 
    presentation.
    
    Parameters:
    data (pd.DataFrame): The input data frame containing the data to be plotted.
    categorical_min (int, optional): The minimum number of unique values a column must have to be considered for 
                                     plotting. Defaults to 2.
    continuous_max (int, optional): The maximum value a column can have to be considered for plotting. Defaults to 100.
    save_to_ppt (bool, optional): Whether to save the plots to a PowerPoint presentation instead of displaying them. 
                                  Defaults to False.
    ppt_filename (str, optional): The filename to use for saving the PowerPoint presentation if save_to_ppt is True. 
                                  Defaults to 'output.pptx'.
    plots_at_a_time (int or 'all', optional): The number of plots to display at a time if save_to_ppt is False, or 
                                              'all' to display all plots at once. Defaults to 'all'.
    
    Raises:
    InvalidInputError: If any of the input parameters do not meet their respective conditions.
    
    Returns:
    None
    """
    try:
        validate_inputs(data, categorical_min, continuous_max, plots_at_a_time)
    except InvalidInputError as e:
        print(e)
        return

    cols_to_plot = get_cols_to_plot(data, categorical_min, continuous_max)

    if not cols_to_plot:
        print("No columns to plot.")
        return

    if save_to_ppt:
        prs = Presentation()

    column_combinations = list(combinations(cols_to_plot, 2))
    num_plots = len(column_combinations)
    if plots_at_a_time == 'all':
        plots_at_a_time = num_plots

    for i, (col1, col2) in enumerate(column_combinations):
        x = data[col1]
        y = data[col2]

        fig, ax = plt.subplots()

        sns.boxplot(x=x, y=y, ax=ax)
        set_plot_properties(ax, col1, col2)

        img_stream = BytesIO()
        plt.tight_layout()
        plt.savefig(img_stream, format='png')
        if save_to_ppt:
            save_plot_to_ppt(prs, img_stream)
            plt.close()
        else:
            plt.show()

        if (i+1) % plots_at_a_time == 0 and i < num_plots - 1:
            input("Press Enter to see next set of plots...")

    if save_to_ppt:
        prs.save(ppt_filename)
        print(f'Saved PowerPoint file to: {ppt_filename}')


# Usage example
np.random.seed(123)
data = {
    'var1': np.random.choice(list(range(10)), 100),
    'var2': np.random.choice(list(range(10)), 100),
    'var3': np.random.choice(list(range(10)), 100),
    'var4': np.random.choice(list(range(20)), 100),
    'var5': np.random.choice(list(range(20)), 100),
    'var6': np.random.choice(list(range(20)), 100),
    'var7': np.random.choice(list(range(200)), 100),
    'var8': np.random.choice(list(range(200)), 100),
    'var9': np.random.choice(list(range(200)), 100),
    'var10': np.random.choice(list(range(200)), 100)
}
df = pd.DataFrame(data)
all_boxes(df, categorical_min=2, continuous_max=200, save_to_ppt=False, plots_at_a_time=4)