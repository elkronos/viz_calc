import numpy as np
import pandas as pd
from pptx import Presentation
from io import BytesIO
import matplotlib.pyplot as plt
from pptx.util import Inches

def all_bars(data: pd.core.frame.DataFrame, unique_levels: int = 10, save_to_ppt: bool = False, add_labels: bool = False, ppt_filename: str = 'output.pptx', bar_color: str = 'skyblue', plots_at_a_time: int = 'all', label_font_size: int = 10) -> None:
    """
    Plot bar charts for columns in a DataFrame with a specified number of unique levels or fewer.

    :param data: Input DataFrame
    :param unique_levels: Maximum number of unique levels a column can have to be plotted (default: 10)
    :param save_to_ppt: Boolean indicating whether to save the plots to a PowerPoint file (default: False)
    :param add_labels: Boolean indicating whether to add labels to the bars (default: False)
    :param ppt_filename: Name of the PowerPoint file to save (if save_to_ppt is True) (default: 'output.pptx')
    :param bar_color: Color of the bars in the plot (default: 'skyblue')
    :param plots_at_a_time: Number of plots to show at a time, use 'all' to show all plots at once (default: 'all')
    :param label_font_size: Font size of the labels on the bars (default: 10)
    """
    def set_plot_properties(ax, col):
        ax.set_ylabel('Count')
        ax.set_xlabel(col)
        ax.set_title(f'{col} Value Counts')

    if not isinstance(data, pd.core.frame.DataFrame):
        print("The input data is not a valid pandas DataFrame.")
        return

    cols_to_plot = [col for col in data.columns if (nunique := data[col].nunique()) <= unique_levels]

    if not cols_to_plot:
        print("No columns to plot.")
        return

    if save_to_ppt:
        prs = Presentation()

    num_plots = len(cols_to_plot)
    if plots_at_a_time == 'all':
        plots_at_a_time = num_plots

    for i, col in enumerate(cols_to_plot, start=1):
        counts = data[col].value_counts().sort_index()
        ax = counts.plot(kind='bar', color=bar_color)
        set_plot_properties(ax, col)
        
        if add_labels:
            for j, v in enumerate(counts):
                ax.text(j, v + 0.5, str(v), ha='center', va='bottom', fontsize=label_font_size)

        if save_to_ppt:
            img_stream = BytesIO()
            plt.tight_layout()
            plt.savefig(img_stream, format='png')
            plt.close()
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            left = top = Inches(1)
            slide.shapes.add_picture(img_stream, left, top, Inches(4.5), Inches(4.5))
        else:
            plt.tight_layout()
            plt.show()
        
        if i % plots_at_a_time == 0 and i < num_plots:
            input("Press Enter to see next set of plots...")
            
    if save_to_ppt:
        prs.save(ppt_filename)
        print(f'Saved PowerPoint file to: {ppt_filename}')


# Usage example:
np.random.seed(123)

data = {}
for i in range(65, 75):
    data[chr(i)] = np.random.choice(list(range(10)), 1000)

for i in range(75, 85):
    data[chr(i)] = np.random.choice(list(range(20)), 1000)

df = pd.DataFrame(data)

all_bars(df, unique_levels=10, save_to_ppt=False, add_labels=True, plots_at_a_time=4)
